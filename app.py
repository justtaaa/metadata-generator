import os
import io
import json
import zipfile
from flask import Flask, request, render_template, jsonify, send_file, send_from_directory
from pptx import Presentation
import fitz  # PyMuPDF
from pydantic import BaseModel, Field
from typing import List, Literal
from datetime import datetime
import shutil
import pandas as pd
import spacy
from sklearn.feature_extraction.text import TfidfVectorizer
from gensim import corpora, models
import re
from PIL import Image
import exifread
import pytesseract
from werkzeug.utils import secure_filename
from openpyxl import load_workbook
import uuid

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads/'
app.config['JSON_FOLDER'] = 'json_outputs/'  # Directory to save JSON files

# Load NLP model
nlp = spacy.load('en_core_web_sm')

# Define controlled vocabularies
SubjectASPVocab = Literal[
    "Primary computing education", "Primary STEM education",
    "Elementary school computing education", "Elementary school STEM education",
    "Middle school computing education", "Middle school STEM education",
    "Secondary computing education", "Secondary STEM education",
    "High school computing education", "High school STEM education",
    "K-12 computing education", "K12/K-12 STEM education",
    "Computer Science", "Python", "MicroPython", "Computer Engineering",
    "Robotics", "Internet of Things (IoT)", "Machine learning (ML)",
    "Artificial intelligence (AI)", "Teach with physical computing",
    "micro:bit", "micro:bit v1", "micro:bit v2",
    "Raspberry Pi", "Raspberry Pi Pico", "Arduino",
    "Computing", "Coding", "Data Science"
]

SubjectAUPVocab = Literal[
    "Computer Science", "Computer Engineering", "Electrical Engineering",
    "Robotics", "Internet of Things (IoT)", "Machine learning (ML)",
    "Artificial intelligence (AI)", "Embedded Systems",
    "Real Time Operating Systems (RTOS)", "Mobile Computing",
    "Cloud Computing", "Edge Computing", "SW Design & Development",
    "Digital System", "Digital Signal Processing", "System-on-Chip Design",
    "Computer Architecture", "VLSI", "Operating Systems", "Linux",
    "MVE / Helium", "Computing"
]

TypeVocab = Literal[
    "EdKit", "Lecture", "Lab", "Video", "Animation", "Course", "Resource"
]

FormatVocab = Literal["ppt", "doc", "zip", "mp3", "pdf", "xlsx", "jpg", "jpeg", "png"]

class FileMetadata(BaseModel):
    title: str = Field(description="The name given to the resource by the creator or publisher")
    creator: str = Field(description="The person or organization primarily responsible for the intellectual content of the resource")
    subject_asp: List[SubjectASPVocab] = Field(description="The Arm School Program subject of the resource")
    subject_aup: List[SubjectAUPVocab] = Field(description="The Arm University Program subject of the resource")
    description: str = Field(description="A textual description of the content of the resource")
    publisher: str = Field(description="The entity responsible for making the resource available")
    contributor: str = Field(description="A person or organization (other than the Creator) who is responsible for making significant contributions to the intellectual content of the resource")
    date: str = Field(description="A date associated with the creation or availability of the resource")
    type: List[TypeVocab] = Field(description="The nature or genre of the content of the resource")
    format: List[FormatVocab] = Field(description="The physical or digital manifestation of the resource")
    identifier: str = Field(description="An unambiguous reference that uniquely identifies the resource within a given context")
    source: str = Field(description="A reference to a second resource from which the present resource is derived")
    language: str = Field(description="The language of the intellectual content of the resource")
    relation: str = Field(description="A reference to a related resource, and the nature of its relationship")
    keywords: List[str] = Field(description="Keywords used")

# Function to clean up JSON output directory
def clean_json_folder():
    folder = app.config['JSON_FOLDER']
    for filename in os.listdir(folder):
        file_path = os.path.join(folder, filename)
        try:
            if os.path.isfile(file_path) or os.path.islink(file_path):
                os.unlink(file_path)
        except Exception as e:
            print(f'Failed to delete {file_path}. Reason: {e}')

# Functions to extract text
def extract_text_from_txt(file_path):
    with open(file_path, 'r', errors='ignore') as f:
        return f.read()

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = []
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        text.append(page.get_text())
    return "\n".join(text)

def extract_text_from_excel(file_path):
    # Using pandas to read Excel file content
    df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets into a dictionary
    content = []
    for sheet_name, sheet_data in df.items():
        content.append(f"Sheet: {sheet_name}")
        
        # Ensure all column names are strings
        sheet_data.columns = sheet_data.columns.astype(str)
        
        # Drop 'Unnamed' columns and NaN values
        sheet_data = sheet_data.dropna(axis=1, how='all')  # Drop columns with all NaN values
        sheet_data = sheet_data.dropna(axis=0, how='all')  # Drop rows with all NaN values
        
        # Check again to ensure all column names are treated as strings
        if sheet_data.columns.str.contains('^Unnamed').any():
            sheet_data = sheet_data.loc[:, ~sheet_data.columns.str.contains('^Unnamed')]
        
        # Convert DataFrame to string
        content.append(sheet_data.to_string(index=False))
    
    # Remove extra spaces and lines
    cleaned_content = "\n".join(content).replace('NaN', '').strip()
    return cleaned_content


# Function to extract text from image files using OCR
def extract_text_from_image(file_path):
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image)
    return text

# Function to extract metadata from image files 
def extract_metadata_from_image(file_path):
    image = Image.open(file_path)
    metadata = image._getexif()  # Extract EXIF data
    if metadata:
        return metadata
    return {}


# Function to clean description
def clean_description(text):
    # Remove excessive newlines and whitespace
    text = re.sub(r'\s+', ' ', text)
    # Remove any non-printable characters or unnecessary symbols
    text = re.sub(r'[^\x00-\x7F]+', ' ', text)
    # Further cleaning logic can be added here if needed
    return text.strip()


# Advanced Metadata Extraction using NLP
def extract_keywords(content):
    tfidf = TfidfVectorizer(max_features=10, stop_words='english')
    tfidf_matrix = tfidf.fit_transform([content])
    keywords = tfidf.get_feature_names_out()
    return list(keywords)

def extract_topics(content):
    tokens = [word.text for word in nlp(content) if not word.is_stop and not word.is_punct]
    dictionary = corpora.Dictionary([tokens])
    corpus = [dictionary.doc2bow(tokens)]
    lda_model = models.LdaModel(corpus, num_topics=1, id2word=dictionary, passes=10)
    topics = lda_model.show_topics(num_words=3)
    return topics

def extract_summary(content):
    doc = nlp(content)
    sentences = list(doc.sents)
    summary = " ".join([sent.text for sent in sentences[:2]])
    return clean_description(summary)  # Apply cleaning here

def map_to_subject_asp(topics):
    mapping = {
        "computing": "Computer Science",
        "python": "Python",
        "microbit": "micro:bit",
        "robotics": "Robotics",
        "machine learning": "Machine learning (ML)",
        "artificial intelligence": "Artificial intelligence (AI)",
        "stem": "Primary STEM education",
        "education": "K-12 computing education",
        "data science": "Data Science",
        # Add more mappings 
    }

    matched_subjects = []
    for topic in topics:
        for key in mapping.keys():
            if key in topic.lower():
                matched_subjects.append(mapping[key])
                break

    if not matched_subjects:
        matched_subjects = ["Computer Science"]  # Default if no matches are found

    return matched_subjects


def detect_metadata(file_path, file_extension, content=None):
    creator = "Unknown"
    publisher = "Unknown"
    creation_date = None

    try:
        if file_extension == '.pptx':
            presentation = Presentation(file_path)
            core_props = presentation.core_properties
            creator = core_props.author or core_props.last_modified_by or "Unknown"
            creation_date = core_props.created

        elif file_extension == '.pdf':
            doc = fitz.open(file_path)
            metadata = doc.metadata
            creator = metadata.get("author") or metadata.get("creator") or "Unknown"
            creation_date = metadata.get("creationDate")
            if creation_date:
                creation_date = datetime.strptime(creation_date, "D:%Y%m%d%H%M%S").strftime("%Y-%m-%d")

        elif file_extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif']:
            with open(file_path, 'rb') as image_file:
                tags = exifread.process_file(image_file)
                if "Image Artist" in tags:
                    creator = tags["Image Artist"].values
                elif "Image Software" in tags:
                    publisher = tags["Image Software"].values
                if "EXIF DateTimeOriginal" in tags:
                    creation_date = tags["EXIF DateTimeOriginal"].values
                    creation_date = datetime.strptime(creation_date, "%Y:%m:%d %H:%M:%S").strftime("%Y-%m-%d")

        elif file_extension == '.xlsx':
            workbook = load_workbook(file_path, read_only=True)
            properties = workbook.properties
            creator = properties.creator or "Unknown"
            creation_date = properties.created.strftime("%Y-%m-%d") if properties.created else None

        elif file_extension == '.txt' and content:
            match = re.search(r'\bby ([A-Z][a-z]+(?: [A-Z][a-z]+)*)\b', content, re.IGNORECASE)
            if match:
                creator = match.group(1)
            

        # Additional logic for other formats 

    except Exception as e:
        print(f"Error detecting metadata for {file_path}: {e}")

    return creator, publisher, creation_date




# Function to combine text extraction with metadata extraction
def extract_metadata(filename, content, image_metadata=None):
    keywords = extract_keywords(content)
    topics = extract_topics(content)
    summary = extract_summary(content)

    # Map extracted topics to controlled vocabulary
    mapped_subject_asp = map_to_subject_asp([t[1] for t in topics])

    # Format mapping
    format_mapping = {
        '.xlsx': 'xlsx',
        '.pdf': 'pdf',
        '.pptx': 'ppt',
        '.doc': 'doc',
        '.png': 'png',
        '.jpg': 'jpg',
        '.jpeg': 'jpeg',
        '.bmp': 'bmp',
        '.tiff': 'tiff',
        '.gif': 'gif'
    }

    # Extract the file extension and use it to get the corresponding format
    file_extension = os.path.splitext(filename)[1].lower()  # Get the file extension and convert to lowercase
    file_format = format_mapping.get(file_extension, 'unknown')  # Default to 'unknown' if not found

    # Detect creator, publisher, and creation date
    creator, publisher, creation_date = detect_metadata(file_path=filename, file_extension=file_extension, content=content)

    # Convert creation_date to a string if it's not None
    if isinstance(creation_date, datetime):
        creation_date = creation_date.strftime("%Y-%m-%d")
    elif creation_date is None:
        creation_date = datetime.now().strftime("%Y-%m-%d")  # Default to current date if not found

    unique_identifier = str(uuid.uuid4())

    file_metadata = FileMetadata(
        title=filename,
        creator=creator,  
        subject_asp=mapped_subject_asp,  # Mapped subject areas
        subject_aup=[],  # Using empty list for demonstration
        description=clean_description(summary),
        publisher=publisher,  
        contributor="Unknown",
        date=creation_date,
        type=["Resource"],
        format=[file_format],
        identifier=unique_identifier,
        source="Original Source",
        language="English",
        relation="None",
        keywords=keywords
    )
    return file_metadata.dict()

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    clean_json_folder()
    session_id = datetime.now().strftime("%Y%m%d%H%M%S")
    session_dir = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
    os.makedirs(session_dir, exist_ok=True)

    if 'textFile' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['textFile']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    file_path = os.path.join(session_dir, secure_filename(file.filename))
    file.save(file_path)

    if file.filename.endswith('.zip'):
        return process_zip_file(file_path, session_dir)

    content = ""
    image_metadata = None

    if file.filename.endswith('.txt'):
        content = extract_text_from_txt(file_path)
    elif file.filename.endswith('.pptx'):
        content = extract_text_from_pptx(file_path)
    elif file.filename.endswith('.pdf'):
        content = extract_text_from_pdf(file_path)
    elif file.filename.endswith('.xlsx'):
        content = extract_text_from_excel(file_path)
    elif file.filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif')):
        content = extract_text_from_image(file_path)
        image_metadata = extract_metadata_from_image(file_path)
    else:
        return jsonify({"error": "Unsupported file type. Please upload a TXT, PPTX, PDF, XLSX, or ZIP file."}), 400

    metadata = extract_metadata(file.filename, content, image_metadata)

    json_filename = os.path.splitext(file.filename)[0] + '.json'
    json_filepath = os.path.join(app.config['JSON_FOLDER'], json_filename)
    os.makedirs(app.config['JSON_FOLDER'], exist_ok=True)
    with open(json_filepath, 'w') as json_file:
        json.dump(metadata, json_file, indent=4)

    return jsonify({
        "message": "File uploaded successfully and JSON file created",
        "metadata": metadata,
        "json_file_path": json_filepath
    }), 200

def process_zip_file(zip_filename, session_dir):
    extracted_folder = os.path.join(session_dir, 'extracted')
    os.makedirs(extracted_folder, exist_ok=True)

    with zipfile.ZipFile(zip_filename, 'r') as zip_ref:
        zip_ref.extractall(extracted_folder)

    processed_files = []
    for root, _, files in os.walk(extracted_folder):
        for file in files:
            file_path = os.path.join(root, file)

            if file.startswith('.') or '__MACOSX' in file_path:
                continue

            content = ""
            image_metadata = None

            if file.endswith('.txt'):
                content = extract_text_from_txt(file_path)
            elif file.endswith('.pptx'):
                content = extract_text_from_pptx(file_path)
            elif file.endswith('.pdf'):
                content = extract_text_from_pdf(file_path)
            elif file.endswith('.xlsx'):
                content = extract_text_from_excel(file_path)
            elif file.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif')):
                content = extract_text_from_image(file_path)
                image_metadata = extract_metadata_from_image(file_path)
            else:
                continue

            # Pass the full file path to the detect_creator function
            metadata = extract_metadata(file_path, content, image_metadata)

            json_filename = os.path.splitext(file)[0] + '.json'
            json_filepath = os.path.join(app.config['JSON_FOLDER'], json_filename)
            with open(json_filepath, 'w') as json_file:
                json.dump(metadata, json_file, indent=4)

            processed_files.append({
                "filename": file,
                "json_file_path": json_filepath
            })

    shutil.rmtree(session_dir)

    return jsonify({
        "message": f"Processed {len(processed_files)} files from ZIP archive",
        "processed_files": processed_files
    }), 200

@app.route('/update_json', methods=['POST'])
def update_json():
    try:
        data = request.json
        json_file_path = data['json_file_path']
        modified_metadata = data['modified_metadata']

        if not os.path.exists(json_file_path):
            return jsonify({"error": "JSON file not found."}), 404

        with open(json_file_path, 'w') as json_file:
            json.dump(modified_metadata, json_file, indent=4)

        return jsonify({"message": "JSON file updated successfully."}), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/json_outputs/<path:filename>', methods=['GET'])
def get_json_file(filename):
    json_directory = app.config['JSON_FOLDER']
    try:
        return send_from_directory(json_directory, filename)
    except FileNotFoundError:
        return jsonify({"error": "JSON file not found."}), 404

@app.route('/download_json', methods=['GET'])
def download_json():
    json_file_path = request.args.get('json_file_path')
    if not json_file_path or not os.path.exists(json_file_path):
        return jsonify({"error": "JSON file not found."}), 404
    
    return send_file(json_file_path, as_attachment=True)


@app.route('/download_all', methods=['GET'])
def download_all():
    json_folder = app.config['JSON_FOLDER']
    
    # Create an in-memory ZIP file
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        # Add all JSON files to the ZIP
        for filename in os.listdir(json_folder):
            if filename.endswith('.json'):
                file_path = os.path.join(json_folder, filename)
                zip_file.write(file_path, arcname=filename)
    
    zip_buffer.seek(0)  # Set the pointer back to the start of the BytesIO object

    # Send the ZIP file for download
    return send_file(zip_buffer, as_attachment=True, download_name='all_metadata.zip', mimetype='application/zip')


if __name__ == '__main__':
    app.run(debug=True)
